# app/services/file_processor.py
"""
Service para procesamiento de diferentes tipos de archivos.
Soporta: PDF, PPTX, PPT, DOCX, DOC, TXT
"""
import io
from typing import Tuple
from fastapi import UploadFile, HTTPException, status
import PyPDF2
import pdfplumber
from pptx import Presentation
from docx import Document
from app.config import settings
from app.core.file_validator import FileValidator


class FileProcessor:
    """Service para procesar y extraer texto de archivos."""

    ALLOWED_EXTENSIONS = settings.ALLOWED_EXTENSIONS
    MAX_FILE_SIZE_BYTES = settings.MAX_FILE_SIZE_MB * 1024 * 1024

    @staticmethod
    def validate_file(file: UploadFile) -> Tuple[str, str]:
        """
        Valida el archivo subido.

        Args:
            file: Archivo subido

        Returns:
            Tupla (nombre_archivo, extensión)

        Raises:
            HTTPException: Si el archivo no es válido
        """
        if not file.filename:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="El archivo no tiene nombre"
            )

        # Obtener extensión
        file_extension = file.filename.split(".")[-1].lower()

        # Validar extensión
        if file_extension not in FileProcessor.ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Formato de archivo no soportado. Formatos permitidos: {', '.join(FileProcessor.ALLOWED_EXTENSIONS)}"
            )

        return file.filename, file_extension

    @staticmethod
    async def validate_file_security(file: UploadFile) -> Tuple[str, str]:
        """
        Valida el archivo usando magic numbers para seguridad.
        
        Args:
            file: Archivo subido
            
        Returns:
            Tupla (nombre_archivo, extensión)
            
        Raises:
            HTTPException: Si el archivo no es válido o es peligroso
        """
        if not file.filename:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="El archivo no tiene nombre"
            )

        # Obtener extensión
        file_extension = file.filename.split(".")[-1].lower()

        # Validar extensión
        if file_extension not in FileProcessor.ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Formato de archivo no soportado. Formatos permitidos: {', '.join(FileProcessor.ALLOWED_EXTENSIONS)}"
            )

        # Validar contenido usando magic numbers
        await FileValidator.validate_file_content(file, file_extension)

        return file.filename, file_extension

    @staticmethod
    async def extract_text(file: UploadFile) -> str:
        """
        Extrae el texto de un archivo según su tipo.

        Args:
            file: Archivo subido

        Returns:
            Texto extraído del archivo

        Raises:
            HTTPException: Si hay error al procesar el archivo
        """
        filename, extension = FileProcessor.validate_file(file)

        # Leer contenido del archivo
        content = await file.read()

        # Validar tamaño
        if len(content) > FileProcessor.MAX_FILE_SIZE_BYTES:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"El archivo excede el tamaño máximo de {settings.MAX_FILE_SIZE_MB}MB"
            )

        # Extraer texto según el tipo de archivo
        try:
            if extension == "txt":
                text = FileProcessor._extract_from_txt(content)
            elif extension == "pdf":
                text = FileProcessor._extract_from_pdf(content)
            elif extension in ["pptx", "ppt"]:
                text = FileProcessor._extract_from_pptx(content)
            elif extension in ["docx", "doc"]:
                text = FileProcessor._extract_from_docx(content)
            else:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Formato no soportado: {extension}"
                )

            # Validar que se extrajo texto
            if not text or len(text.strip()) < 10:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="No se pudo extraer texto del archivo o el contenido es muy corto"
                )

            return text.strip()

        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error al procesar el archivo: {str(e)}"
            )

    @staticmethod
    def _extract_from_txt(content: bytes) -> str:
        """Extrae texto de un archivo TXT."""
        try:
            # Intentar UTF-8 primero, luego latin-1
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1")
        except Exception as e:
            raise Exception(f"Error al leer archivo TXT: {str(e)}")

    @staticmethod
    def _extract_from_pdf(content: bytes) -> str:
        """Extrae texto de un archivo PDF usando pdfplumber (más robusto)."""
        text_parts = []

        try:
            # Usar pdfplumber para mejor extracción de texto
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            # Si pdfplumber no funciona, intentar con PyPDF2
            if not text_parts:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            return "\n\n".join(text_parts)

        except Exception as e:
            raise Exception(f"Error al leer PDF: {str(e)}")

    @staticmethod
    def _extract_from_pptx(content: bytes) -> str:
        """Extrae texto de un archivo PPTX."""
        text_parts = []

        try:
            prs = Presentation(io.BytesIO(content))

            for slide in prs.slides:
                # Extraer texto de todas las formas en la diapositiva
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)

                    # Si la forma tiene una tabla
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text)

            return "\n\n".join(text_parts)

        except Exception as e:
            raise Exception(f"Error al leer PPTX: {str(e)}")

    @staticmethod
    def _extract_from_docx(content: bytes) -> str:
        """Extrae texto de un archivo DOCX."""
        text_parts = []

        try:
            doc = Document(io.BytesIO(content))

            # Extraer párrafos
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extraer tablas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)

            return "\n\n".join(text_parts)

        except Exception as e:
            raise Exception(f"Error al leer DOCX: {str(e)}")
